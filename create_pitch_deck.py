#!/usr/bin/env python3
"""
CeloQuest Pitch Deck Generator
Creates a professional PowerPoint presentation for CeloQuest
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide():
    """Slide 1: Title Slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(252, 211, 77)  # Yellow
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "CeloQuest"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Gamified Micro-Lending on Celo"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.6))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Fund entrepreneurs worldwide with as little as one dollar"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(20)
    tagline_para.font.color.rgb = RGBColor(53, 208, 127)
    tagline_para.alignment = PP_ALIGN.CENTER

def add_problem_slide():
    """Slide 2: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Problem"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(251, 191, 36)
    
    content = slide.placeholders[1].text_frame
    content.text = "Over 1.7 billion people worldwide lack access to traditional banking"
    
    for bullet_text in [
        "Entrepreneurs in developing countries struggle to secure small business loans",
        "High barriers to entry and lack of credit history",
        "Limited access to traditional financial institutions",
        "People wanting to make social impact lack accessible platforms"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(20)

def add_solution_slide():
    """Slide 3: Our Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Our Solution"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(251, 191, 36)
    
    content = slide.placeholders[1].text_frame
    content.text = "A decentralized micro-lending platform on Celo blockchain"
    
    for bullet_text in [
        "Lend as little as $1 to verified entrepreneurs using stablecoins (cUSD)",
        "Instant, affordable cross-border transactions",
        "Gamified experience with impact points and badges",
        "Complete transparency and portfolio tracking",
        "No intermediaries, direct peer-to-peer lending"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(20)

def add_how_it_works_slide():
    """Slide 4: How It Works"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "How It Works"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(251, 191, 36)
    
    content = slide.placeholders[1].text_frame
    content.text = "1. Browse verified entrepreneurs and their stories"
    
    for step, bullet_text in enumerate([
        "2. Lend any amount starting from $1 using Celo stablecoins",
        "3. Track your portfolio and borrower progress",
        "4. Earn impact points and unlock achievement badges",
        "5. Swap between CELO, cUSD, and cEUR seamlessly"
    ], start=1):
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(20)

def add_market_slide():
    """Slide 5: Market Opportunity"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Market Opportunity"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(251, 191, 36)
    
    content = slide.placeholders[1].text_frame
    content.text = "$380 billion global micro-lending market"
    
    for bullet_text in [
        "1.7 billion unbanked adults worldwide",
        "Growing cryptocurrency adoption in emerging markets",
        "Celo's mobile-first approach perfect for target demographics",
        "96.5% average repayment rate in micro-lending sector"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(20)

def add_technology_slide():
    """Slide 6: Technology Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technology Stack"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(251, 191, 36)
    
    content = slide.placeholders[1].text_frame
    content.text = "Built on Celo Blockchain"
    
    for bullet_text in [
        "Smart contracts for transparent, automated lending",
        "Next.js & React for responsive web interface",
        "ethers.js for blockchain interactions",
        "Ubeswap integration for seamless token swaps",
        "EmailJS for entrepreneur application system"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(20)

def add_traction_slide():
    """Slide 7: Traction & Metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Gamification Features"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(251, 191, 36)
    
    content = slide.placeholders[1].text_frame
    content.text = "Engaging users through game mechanics"
    
    for bullet_text in [
        "Impact Points: Earn points for every dollar lent",
        "Achievement Badges: Bronze, Silver, Gold, Platinum tiers",
        "Portfolio Dashboard: Track your lending history",
        "Leaderboards: Compete with other lenders (coming soon)",
        "NFT Rewards: Unlock special badges as NFTs (planned)"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(20)

def add_roadmap_slide():
    """Slide 8: Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Roadmap"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(251, 191, 36)
    
    content = slide.placeholders[1].text_frame
    content.text = "Q1 2026: Launch MVP with verified entrepreneurs"
    
    for bullet_text in [
        "Q2 2026: Implement repayment tracking & NFT badges",
        "Q3 2026: Mobile app launch (iOS & Android)",
        "Q4 2026: Partnerships with microfinance organizations",
        "2027: Expand to 10+ countries, DAOize governance"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(20)

def add_team_slide():
    """Slide 9: Team"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Why Celo?"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(251, 191, 36)
    
    content = slide.placeholders[1].text_frame
    content.text = "Perfect blockchain for financial inclusion"
    
    for bullet_text in [
        "Mobile-first design for smartphone accessibility",
        "Ultra-low transaction fees (~$0.01)",
        "Carbon-negative blockchain (environmentally friendly)",
        "Built-in stablecoins (cUSD, cEUR) for stability",
        "Strong focus on real-world use cases and social impact"
    ]:
        p = content.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(20)

def add_contact_slide():
    """Slide 10: Contact"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(252, 211, 77)
    
    # Thank you text
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Join Us in Empowering\nEntrepreneurs Worldwide"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_frame.text = "GitHub: github.com/AdekunleBamz/celoquest\n"
    
    for line in [
        "Website: celoquest.vercel.app",
        "Built with ❤️ on Celo"
    ]:
        p = contact_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

# Generate all slides
print("Generating CeloQuest pitch deck...")
add_title_slide()
add_problem_slide()
add_solution_slide()
add_how_it_works_slide()
add_market_slide()
add_technology_slide()
add_traction_slide()
add_roadmap_slide()
add_team_slide()
add_contact_slide()

# Save presentation
output_file = 'CeloQuest-Pitch-Deck.pptx'
prs.save(output_file)
print(f"✅ Pitch deck created successfully: {output_file}")
print(f"📊 Total slides: 10")
